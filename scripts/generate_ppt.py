from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


# 使用新文件名生成修订版，避免 PowerPoint 打开旧文件时发生 Windows 文件锁冲突。
OUT = Path("docs/BeamDrop_答辩汇报_高对比度版.pptx")
W, H = Inches(13.333), Inches(7.5)
BG = RGBColor(11, 18, 32)
PANEL = RGBColor(20, 32, 51)
PANEL_2 = RGBColor(28, 43, 66)
ACCENT = RGBColor(32, 199, 217)
ACCENT_2 = RGBColor(80, 225, 190)
WHITE = RGBColor(255, 255, 255)
# 黑色背景上不使用灰色正文；此颜色仅作低层级说明，仍保持高对比度。
MUTED = RGBColor(220, 235, 248)
RED = RGBColor(255, 112, 112)


def set_bg(slide):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = BG
    # 顶部点缀线
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.06))
    line.fill.solid(); line.fill.fore_color.rgb = ACCENT; line.line.fill.background()


def textbox(slide, text, x, y, w, h, size=20, color=WHITE, bold=False,
            align=PP_ALIGN.LEFT, font="Microsoft YaHei", valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear(); tf.word_wrap = True; tf.vertical_anchor = valign
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line; p.alignment = align
        p.space_after = Pt(5)
        for r in p.runs:
            r.font.name = font; r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    return box


def title(slide, heading, sub=None):
    textbox(slide, heading, Inches(.65), Inches(.42), Inches(11.9), Inches(.48), 26, WHITE, True)
    if sub:
        textbox(slide, sub, Inches(.68), Inches(.95), Inches(11.6), Inches(.34), 15, WHITE)


def footer(slide, page):
    textbox(slide, "BeamDrop（邻光）｜局域网文件传输系统", Inches(.68), Inches(7.04), Inches(7), Inches(.24), 11, WHITE)
    textbox(slide, f"{page:02d}", Inches(12.2), Inches(7.00), Inches(.45), Inches(.26), 12, ACCENT, True, PP_ALIGN.RIGHT)


def rounded(slide, x, y, w, h, fill=PANEL, radius=True):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill; shp.line.color.rgb = fill
    return shp


def card(slide, x, y, w, h, head, body, accent=ACCENT):
    rounded(slide, x, y, w, h)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(.07), h)
    bar.fill.solid(); bar.fill.fore_color.rgb = accent; bar.line.fill.background()
    textbox(slide, head, x+Inches(.25), y+Inches(.18), w-Inches(.45), Inches(.36), 18, WHITE, True)
    textbox(slide, body, x+Inches(.25), y+Inches(.62), w-Inches(.45), h-Inches(.75), 15, WHITE)


def bullet_list(slide, items, x, y, w, size=18, color=WHITE, gap=.56):
    for idx, item in enumerate(items):
        textbox(slide, "●", x, y+Inches(idx*gap), Inches(.25), Inches(.26), size-4, ACCENT)
        textbox(slide, item, x+Inches(.32), y+Inches(idx*gap), w-Inches(.32), Inches(.40), size, color)


def arrow(slide, x, y, w, h=Inches(.25), color=ACCENT):
    s = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()


def node(slide, x, y, w, h, text, fill=PANEL_2, size=14):
    rounded(slide, x, y, w, h, fill)
    textbox(slide, text, x+Inches(.08), y+Inches(.12), w-Inches(.16), h-Inches(.22), max(size, 13), WHITE, True, PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)


def placeholder(slide, x, y, w, h, label, note):
    shp = rounded(slide, x, y, w, h, PANEL_2)
    shp.line.color.rgb = ACCENT
    textbox(slide, "▧", x, y+Inches(.45), w, Inches(.55), 34, ACCENT, True, PP_ALIGN.CENTER)
    textbox(slide, label, x+Inches(.2), y+Inches(1.15), w-Inches(.4), Inches(.38), 16, WHITE, True, PP_ALIGN.CENTER)
    textbox(slide, note, x+Inches(.28), y+Inches(1.62), w-Inches(.56), h-Inches(1.82), 15, WHITE, False, PP_ALIGN.CENTER)


def main():
    prs = Presentation(); prs.slide_width = W; prs.slide_height = H
    blank = prs.slide_layouts[6]

    # 1 Cover
    s = prs.slides.add_slide(blank); set_bg(s)
    textbox(s, "BeamDrop", Inches(.78), Inches(1.25), Inches(7.7), Inches(.95), 46, WHITE, True)
    textbox(s, "邻光｜局域网文件传输系统的设计与实现", Inches(.84), Inches(2.25), Inches(8), Inches(.48), 22, ACCENT, True)
    textbox(s, "React · FastAPI · pybind11 · C++ · TCP", Inches(.86), Inches(2.92), Inches(7), Inches(.36), 18, WHITE)
    textbox(s, "课程名称：[待填写]\n汇报人：[待填写]\n指导教师：[待填写]    日期：[待填写]", Inches(.86), Inches(4.45), Inches(5.2), Inches(1.1), 16, WHITE)
    # decorative network
    for x, y, r in [(9.4,1.5,.25),(11.2,2.3,.17),(10.1,3.5,.2),(11.8,4.5,.3),(8.8,5.15,.14)]:
        c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(r), Inches(r)); c.fill.solid(); c.fill.fore_color.rgb=ACCENT; c.line.fill.background()
    footer(s, 1)

    # 2 Background
    s=prs.slides.add_slide(blank); set_bg(s); title(s,"项目背景与问题","不仅要“能传”，还要可观察、可验证、可维护")
    card(s, Inches(.72), Inches(1.55), Inches(2.8), Inches(3.8), "传输过程不可见", "用户无法确认传输是否仍在进行，无法观察完成进度。")
    card(s, Inches(3.77), Inches(1.55), Inches(2.8), Inches(3.8), "失败原因不明确", "连接失败后，任务不应假成功或长期停留在运行状态。", ACCENT_2)
    card(s, Inches(6.82), Inches(1.55), Inches(2.8), Inches(3.8), "正确性难确认", "传输结束不等于内容正确，需要可复核的完整性证据。")
    card(s, Inches(9.87), Inches(1.55), Inches(2.8), Inches(3.8), "原生核心难操作", "底层网络能力需要一个清晰、直观的控制入口。", ACCENT_2)
    textbox(s,"目标：构建从 C++ 传输核心到 Web 控制面的可验证闭环", Inches(.75), Inches(5.85), Inches(11.7), Inches(.4), 22, WHITE, True, PP_ALIGN.CENTER); footer(s,2)

    # 3 scope
    s=prs.slides.add_slide(blank); set_bg(s); title(s,"需求与项目范围","明确 MVP 的已交付能力与非目标，避免夸大项目范围")
    card(s, Inches(.72), Inches(1.45), Inches(5.9), Inches(4.8), "✓ 已实现能力", "• 本机路径的单文件发送\n• 接收服务启动、状态展示和停止\n• pending → running → completed / failed\n• REST 快照恢复与 WebSocket 增量事件\n• SHA256 完整性验证与失败语义", ACCENT_2)
    card(s, Inches(6.75), Inches(1.45), Inches(5.9), Inches(4.8), "— 本阶段不包含", "• UDP 自动发现、认证加密、公网部署\n• 多客户端并发接收、任务历史数据库\n• 桌面安装包、托盘、自启动、移动端\n• 立即中断阻塞 socket 的强制取消", RED)
    footer(s,3)

    # 4 architecture
    s=prs.slides.add_slide(blank); set_bg(s); title(s,"总体架构","控制面与真实 C++ TCP 传输核心分层协作")
    labels=["React GUI\ndashboardState", "FastAPI\nTaskManager / ReceiverManager / EventBus", "pybind11\nbeamdrop_native", "C++ Core\napp service / TCP transfer"]
    xs=[.7,3.35,7.0,9.75]; ws=[2.15,3.0,2.05,2.8]
    for i,(lab,x,w) in enumerate(zip(labels,xs,ws)):
        node(s,Inches(x), Inches(2.05), Inches(w), Inches(1.28),lab, PANEL_2, 14)
        if i<3: arrow(s, Inches(x+w+.12), Inches(2.54), Inches(.43))
    card(s, Inches(1.4), Inches(4.45), Inches(4.7), Inches(1.3), "REST /api/snapshot", "权威状态快照：刷新、重连或序列缺口时恢复。")
    card(s, Inches(7.15), Inches(4.45), Inches(4.7), Inches(1.3), "WebSocket", "增量事件与 heartbeat：负责实时更新，不承载权威状态。", ACCENT_2)
    footer(s,4)

    # 5 layering
    s=prs.slides.add_slide(blank); set_bg(s); title(s,"C++ 核心分层设计","协议与网络解耦，传输流程与上层 UI/CLI 解耦")
    layers=[("app","配置、日志与应用编排"),("filesystem","文件读写、路径处理与路径穿越防护"),("transfer","发送、接收、会话、续传与进度"),("protocol","Packet、序列化、解析与 checksum 校验"),("network","TCP 连接、监听、字节读写")]
    for i,(n,d) in enumerate(layers):
        y=1.4+i*.88; width=10.8-i*.6; x=.82+i*.3
        rounded(s, Inches(x), Inches(y), Inches(width), Inches(.64), PANEL_2 if i%2==0 else PANEL)
        textbox(s,n, Inches(x+.25), Inches(y+.16), Inches(1.7), Inches(.3), 15, ACCENT, True)
        textbox(s,d, Inches(x+2.15), Inches(y+.16), Inches(width-2.4), Inches(.3), 16, WHITE)
    textbox(s,"核心原则：每一层只处理自己的问题，降低扩展与测试的耦合成本。", Inches(.85), Inches(6.2), Inches(11), Inches(.35), 19, WHITE, False, PP_ALIGN.CENTER); footer(s,5)

    # 6 flow
    s=prs.slides.add_slide(blank); set_bg(s); title(s,"关键传输与状态流程","一次任务从文件扫描到安全落盘，并具备明确终态")
    textbox(s,"发送链路", Inches(.82), Inches(1.45), Inches(2), Inches(.35), 18, ACCENT, True)
    send=["文件路径","FileScanner","Sender","Protocol","TcpClient","TCP 网络"]
    for i,t in enumerate(send):
        x=.75+i*2.0; node(s,Inches(x),Inches(1.95),Inches(1.55),Inches(.62),t, PANEL_2,11)
        if i<5: arrow(s, Inches(x+1.6), Inches(2.15), Inches(.3), Inches(.2))
    textbox(s,"接收链路", Inches(.82), Inches(3.15), Inches(2), Inches(.35), 18, ACCENT_2, True)
    recv=["TCP 网络","TcpServer","Protocol","Receiver","FileWriter","保存目录"]
    for i,t in enumerate(recv):
        x=.75+i*2.0; node(s,Inches(x),Inches(3.65),Inches(1.55),Inches(.62),t, PANEL,11)
        if i<5: arrow(s, Inches(x+1.6), Inches(3.85), Inches(.3), Inches(.2), ACCENT_2)
    textbox(s,"任务状态机", Inches(.82), Inches(5.12), Inches(2), Inches(.35), 18, ACCENT, True)
    for x,t,col in [(1.3,"pending",PANEL_2),(3.65,"running",PANEL_2),(6.0,"completed",ACCENT_2),(8.35,"failed\nNETWORK_FAILURE",RED)]:
        node(s,Inches(x),Inches(5.62),Inches(1.75),Inches(.63),t,col,11)
    arrow(s,Inches(3.1),Inches(5.83),Inches(.4)); arrow(s,Inches(5.45),Inches(5.83),Inches(.4)); arrow(s,Inches(7.8),Inches(5.83),Inches(.4),Inches(.2),RED); footer(s,6)

    # 7 GUI
    s=prs.slides.add_slide(blank); set_bg(s); title(s,"GUI 与实时状态设计","浏览器 GUI 是控制面，不是模拟界面")
    placeholder(s, Inches(.72), Inches(1.45), Inches(6.25), Inches(4.95), "GUI 截图占位区", "请替换为 1920×1080 或更高分辨率的\nBeamDrop 页面全屏截图。\n建议包含接收服务、发送文件与任务列表。")
    card(s, Inches(7.35), Inches(1.45), Inches(5.3), Inches(1.08), "接收服务管理", "主机、端口、保存目录、启动/停止状态")
    card(s, Inches(7.35), Inches(2.75), Inches(5.3), Inches(1.08), "发送与进度反馈", "绝对路径、目标地址、任务状态、字节数、进度")
    card(s, Inches(7.35), Inches(4.05), Inches(5.3), Inches(1.08), "结构化失败", "例如 NETWORK_FAILURE：明确进入 failed 终态", RED)
    card(s, Inches(7.35), Inches(5.35), Inches(5.3), Inches(1.08), "状态恢复", "刷新或 WebSocket 重连后由快照恢复", ACCENT_2)
    footer(s,7)

    # 8 demo
    s=prs.slides.add_slide(blank); set_bg(s); title(s,"现场演示：完整传输闭环","建议 3 分钟；普通流程与故障流程均使用预设环境")
    steps=[("01","启动接收服务","127.0.0.1:9090，确认 running"),("02","正常传输","观察 pending → running → completed"),("03","边界文件","空文件与中文文件名文件"),("04","完整性校验","源文件与接收文件 SHA256 一致"),("05","失败语义","目标 127.0.0.1:19091 → failed"),("06","停止服务","确认 native 最终状态为 stopped")]
    for i,(n,h,b) in enumerate(steps):
        row=i//3; col=i%3; x=.75+col*4.15; y=1.45+row*2.2
        card(s,Inches(x),Inches(y),Inches(3.75),Inches(1.65),f"{n}  {h}",b, ACCENT if i not in (3,4) else (ACCENT_2 if i==3 else RED))
    textbox(s,"备用策略：提前录制 2–3 分钟演示视频；环境异常时直接切换录屏。", Inches(.8), Inches(6.3), Inches(11.7), Inches(.3), 17, WHITE, False, PP_ALIGN.CENTER); footer(s,8)

    # 9 evidence
    s=prs.slides.add_slide(blank); set_bg(s); title(s,"测试与验证证据","自动化测试 + 文件哈希 + Linux 端到端场景")
    metrics=[("3 / 3","Windows C++ CTest","核心传输测试通过"),("31 / 31","Python pytest","后端测试通过"),("10 / 10","React test","前端测试通过；lint、build 通过"),("E2E","Linux Ubuntu VM","普通/空/中文文件名、SHA256、失败、恢复、停止")]
    for i,(num,h,b) in enumerate(metrics):
        col=i%2; row=i//2; x=.85+col*6.1; y=1.55+row*2.25
        rounded(s, Inches(x), Inches(y), Inches(5.6), Inches(1.8), PANEL)
        textbox(s,num, Inches(x+.28), Inches(y+.28), Inches(1.55), Inches(.6), 28, ACCENT, True)
        textbox(s,h, Inches(x+1.95), Inches(y+.3), Inches(3.3), Inches(.34), 16, WHITE, True)
        textbox(s,b, Inches(x+1.95), Inches(y+.82), Inches(3.3), Inches(.6), 14, WHITE)
    textbox(s,"说明：Linux CTest 为 2/2；pybind build/import 为独立验证项，不能混写成 Linux CTest 3/3。", Inches(.9), Inches(6.35), Inches(11.4), Inches(.3), 14, WHITE, False, PP_ALIGN.CENTER); footer(s,9)

    # 10 results
    s=prs.slides.add_slide(blank); set_bg(s); title(s,"项目成果","完成了从传输核心到可视化控制面的可验证闭环")
    card(s, Inches(.85), Inches(1.65), Inches(3.65), Inches(3.85), "01  分层可维护", "C++ TCP 核心与协议、网络、文件系统解耦；上层控制面不侵入传输实现。")
    card(s, Inches(4.84), Inches(1.65), Inches(3.65), Inches(3.85), "02  过程可观察", "真实任务的实时进度、结构化失败、接收服务状态与快照恢复。", ACCENT_2)
    card(s, Inches(8.83), Inches(1.65), Inches(3.65), Inches(3.85), "03  结果可验证", "自动化测试、SHA256 完整性校验与跨 Windows/Linux 端到端验证。")
    textbox(s,"真实 C++ 引擎，而非 GUI 模拟数据。", Inches(.85), Inches(6.02), Inches(11.6), Inches(.4), 21, WHITE, True, PP_ALIGN.CENTER); footer(s,10)

    # 11 future
    s=prs.slides.add_slide(blank); set_bg(s); title(s,"限制与后续规划","明确当前边界，按传输能力、体验与安全能力演进")
    card(s, Inches(.75), Inches(1.5), Inches(5.75), Inches(4.8), "当前限制", "• 路径式单文件发送；接收端单实例、串行接收\n• 无自动发现、并发接收、任务持久化\n• 取消是协作式请求，不保证立即打断 socket\n• 默认仅绑定本机回环地址，不是公网部署方案", RED)
    textbox(s,"演进路径", Inches(7.3), Inches(1.55), Inches(3), Inches(.35), 20, ACCENT, True)
    future=["UDP 自动发现","多文件/目录与并发","加密认证","桌面端封装","任务历史与可观测性"]
    for i,t in enumerate(future):
        y=2.1+i*.74; node(s,Inches(7.35),Inches(y),Inches(4.5),Inches(.5),t,PANEL_2,13)
        if i<4: arrow(s,Inches(9.38),Inches(y+.52),Inches(.3),Inches(.15))
    footer(s,11)

    # 12 conclusion
    s=prs.slides.add_slide(blank); set_bg(s)
    textbox(s,"BeamDrop", Inches(.8), Inches(1.35), Inches(5), Inches(.75), 42, WHITE, True)
    textbox(s,"让局域网文件传输具备可维护的核心、\n可观察的过程与可验证的结果。", Inches(.84), Inches(2.25), Inches(7.3), Inches(1.05), 25, ACCENT, True)
    bullet_list(s,["分层架构：控制面与 C++ 传输核心解耦","可观察性：实时进度、确定状态、失败语义、刷新恢复","可验证性：自动测试 + SHA256 + 端到端场景"], Inches(.9), Inches(4.15), Inches(8.6), 18, WHITE,.63)
    textbox(s,"谢谢，欢迎提问", Inches(.88), Inches(6.08), Inches(6), Inches(.4), 22, WHITE, True)
    footer(s,12)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"Created: {OUT.resolve()} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()